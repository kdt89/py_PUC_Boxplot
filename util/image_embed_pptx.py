from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from typing import List
from controller.setting import Setting


class ImageEmbedPPTX:

    PPTX_TEMPLATE = 'Template_Report.pptx'
    PPTX_OUTPUT_NAME = 'Report'
    FIXED_TEXT_BOX_SHAPE_HEIGHT = 0.5
    LIST_IMAGE_FILEPATHS: List[str] = Setting.LIST_FIGURE_IMAGES

    """
    GRAPH AREA IN SLIDE: area to attach external image to it.
    The attached image will be centered in the area (horizontally and vertically)
    Below positioning is based on the Top, Left corner of the slide as (0, 0) coordinate
    """
    GRAPH_TOP_MARGIN = Inches(1.75) # Distance from top edge of Slide to Graph area
    GRAPH_BOT_MARGIN = Inches(0.25) # Distance from bottom edge of Slide to Graph area
    GRAPH_LEFT_MARGIN = Inches(0.15) # Distance from left edge of Slide to Graph area
    GRAPH_RIGHT_MARGIN = Inches(0.15) # Distance from right edge of Slide to Graph area

    SLIDE_WIDTH = 0 # Initialize only, change value after loading Template file
    SLIDE_HEIGHT = 0 # Initialize only, change value after loading Template file
    GRAPH_WIDTH = 0 # Initialize only, change value after loading Template file
    GRAPH_HEIGHT = 0 # Initialize only, change value after loading Template file


    def __init__(self):
        # check file existence
        self.pptx: Presentation = None

        if not os.path.exists(self.PPTX_TEMPLATE):
            # print(f'File {self.PPTX_FILENAME} not found.')
            return
        try:
            self.pptx = Presentation(self.PPTX_TEMPLATE)
        except:
            print(f'Failed to open Template file {self.PPTX_TEMPLATE} for exporting image to.')
            self.pptx = None
            return

        # self.slide_width = self.pptx.slide_width
        # self.slide_height = self.pptx.slide_height
        ImageEmbedPPTX.SLIDE_HEIGHT = self.pptx.slide_height
        ImageEmbedPPTX.SLIDE_WIDTH = self.pptx.slide_width
        ImageEmbedPPTX.GRAPH_HEIGHT = ImageEmbedPPTX.SLIDE_HEIGHT - ImageEmbedPPTX.GRAPH_TOP_MARGIN - ImageEmbedPPTX.GRAPH_BOT_MARGIN
        ImageEmbedPPTX.GRAPH_WIDTH = ImageEmbedPPTX.SLIDE_WIDTH - ImageEmbedPPTX.GRAPH_LEFT_MARGIN - ImageEmbedPPTX.GRAPH_RIGHT_MARGIN


    def clearAllShapes(self) -> None:
        if self.pptx is None:
            return
        
        for slide in self.pptx.slides:
            # Remove the old figures
            shapes = slide.shapes
            for shape in shapes:
                match shape.shape_type:
                    case MSO_SHAPE_TYPE.GROUP | MSO_SHAPE_TYPE.PICTURE | MSO_SHAPE_TYPE.TABLE:
                        shapes.element.remove(shape.element)

    @staticmethod
    def getImageScaleVal(
            from_height: int,
            from_width: int,
            to_height: int,
            to_width: int) -> float:
        # scale figure size to max fit with max_width and max_height
        width_scale = to_width / from_width
        height_scale = to_height / from_height
        final_scale = min(width_scale, height_scale)
        return final_scale


    def exportImages2PPTX(self) -> None:
        if len(self.LIST_IMAGE_FILEPATHS) == 0:
            return

        if self.pptx is None:
            return

        i = 0
        for img in self.LIST_IMAGE_FILEPATHS:
            if i >= len(self.pptx.slides):
                self.pptx.slides.add_slide(self.pptx.slide_layouts[0])
                
            slide_current = self.pptx.slides[i]
            # plot_img = slide_current.shapes.add_picture(img, Inches(0.40), Inches(4.85), width=Inches(5.30))
            added_img = slide_current.shapes.add_picture(image_file=img, left=0, top=0)

            # resize the added image to fit Graph area
            img_scale = ImageEmbedPPTX.getImageScaleVal(
                from_height=int(added_img.height),
                from_width=int(added_img.width),
                to_height=int(ImageEmbedPPTX.GRAPH_HEIGHT),
                to_width=int(ImageEmbedPPTX.GRAPH_WIDTH))
            
            added_img.width = int(added_img.width * img_scale)
            added_img.height = int(added_img.height * img_scale)

            # align picture to the center horizontally and vertically
            added_img.left = int(ImageEmbedPPTX.GRAPH_LEFT_MARGIN + (ImageEmbedPPTX.GRAPH_WIDTH - added_img.width) / 2)
            added_img.top = int(ImageEmbedPPTX.GRAPH_TOP_MARGIN + (ImageEmbedPPTX.GRAPH_HEIGHT - added_img.height) / 2)

            # Send the figures to the back
            ref_element = slide_current.shapes[0]._element
            ref_element.addprevious(added_img._element)
            i += 1

        # Save the PPT
        pptx_output_file = Setting.OUTPUT_DIR + '\\' + self.PPTX_OUTPUT_NAME + '.pptx'
        self.pptx.save(pptx_output_file)